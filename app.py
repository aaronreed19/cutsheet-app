import os
from flask import Flask, render_template, request, redirect, url_for, abort, \
    send_from_directory
from werkzeug.utils import secure_filename
from pptx import Presentation
from PyPDF2 import PdfFileReader
from pptx.util import Inches
from wand.image import Image as wi

app = Flask(__name__)
#app.config['MAX_CONTENT_LENGTH'] = 2 * 1024 * 1024
app.config['UPLOAD_EXTENSIONS'] = ['.pdf','.pptx']
app.config['UPLOAD_PATH'] = 'uploads'

@app.errorhandler(413)
def too_large(e):
    return "File is too large", 413

@app.route('/')
def index():
    dir = os.path.abspath(os.path.join(os.path.dirname( __file__ ), 'uploads'))
    for f in os.listdir(dir):
        os.remove(os.path.join(dir, f))
    files = os.listdir(dir)
    return render_template('index.html', files=files)

@app.route('/', methods=['POST'])
def upload_files():
    dir = os.path.abspath(os.path.join(os.path.dirname( __file__ ), 'uploads'))
    uploaded_file = request.files['file']
    for uploaded_file in request.files.getlist('file'):
        filename = secure_filename(uploaded_file.filename)
        if filename != '':
            file_ext = os.path.splitext(filename)[1]
            if file_ext not in app.config['UPLOAD_EXTENSIONS']:
                return "Invalid image", 400
            uploaded_file.save(os.path.join(dir, filename))

    return '', 204




@app.route('/download', methods=['GET'])
def create_ppt():
    SLD_TITLE = 0   #This is just which template it is aka the first one
    NRML_SLD = 1    
    FixType = 10
    PageNum = 11
    width = None
    height = Inches(8.25)
    left = Inches(.75)
    top = Inches(1.75)


    current_dir = os.path.abspath(os.path.join(os.path.dirname( __file__ ), 'uploads'))   #Just for finding the template
    prs = Presentation(current_dir + '\\actual_template.pptx')

    title_slide_layout = prs.slide_layouts[SLD_TITLE]
    normal_slide = prs.slide_layouts[NRML_SLD]
    title_slide = prs.slides.add_slide(title_slide_layout)  #These two add new slides (one below)


    i = 1
    os.chdir(current_dir)
    # directory = current_dir
    for file_name in os.listdir(current_dir):
        if file_name.endswith(".pdf"):
            try:
                j = 1
                #file_location = directory + '\\' + file_name
                pdf = PdfFileReader(file_name)
                number_of_pages = pdf.getNumPages()
                pdf = wi(filename=file_name, resolution=300)
                pdfimage = pdf.convert("jpeg")
                type = file_name.replace('.pdf','')
                if type[1]== '0':
                    type=type.replace('0','')
                for img in pdfimage.sequence:
                        page = wi(image=img)
        #If page number begins with 0 take out 0
                        page.save(filename=type + '_' + str(i)+".jpg")
                        slide = prs.slides.add_slide(normal_slide)
                        slide.placeholders[FixType].text = type
                        slide.placeholders[PageNum].text = str(j) + ' of ' + str(number_of_pages)
                        shapes = prs.slides[i].shapes
                        shapes.add_picture(type + '_' + str(i)+".jpg", left, top, width, height)
                        os.remove(type + '_' + str(i)+".jpg")
                        i +=1
                        j+=1
            except:
                print('Error: '+file_name+ ' pdf type not supported')



    #os.chdir(current_dir)
    prs.save('test_add.pptx')
    return render_template('home.html')


@app.route("/get-ppt/<ppt_name>")
def get_ppt(ppt_name):
    current_dir = os.path.abspath(os.path.join(os.path.dirname( __file__ ), 'uploads'))
    try:
        return send_from_directory(current_dir, filename=ppt_name, as_attachment=True)
    except FileNotFoundError:
        abort(404)


if __name__ == '__main__':
   app.run()